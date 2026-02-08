
import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Helper function to add slide with title and content
    def add_slide(title, content):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title_placeholder.text = title
        content_placeholder.text = content

    # Helper function to add slide with title and bullets
    def add_bullet_slide(title, bullets):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title_placeholder.text = title
        tf = content_placeholder.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = bullet
            p.level = 0

    # Slide 1: Overview
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Unified Microservices Platform (UMP)"
    subtitle.text = "A Multi-Tenant App Factory\nPowering FitIt, Wissler, and Future Apps"

    # Slide 2: The Core Services
    add_bullet_slide("The Core Services (The 'Kernel')", [
        "Infrastructure: 25+ Dockerized Containers providing day-1 readiness.",
        "Identity Provider (Auth API): Centralized JWT, Social Login (Google, Microsoft, Facebook, Apple).",
        "Payment Engine: Unified multi-tenant gateway (Stripe, PayPal).",
        "Notification Hub: SignalR, SMS, Email dispatcher.",
        "File System (Media API): S3-compatible storage.",
        "Search & Geo: Global text search and location indexing."
    ])

    # Slide 3: Global Config Strategy
    add_bullet_slide("Global Configuration Strategy", [
        "SMS/SMTP: Single Global Pipeline (High Reputation Sender).",
        "Social Login: Multi-Tenant Config (Toggle per App).",
        "Payments: Multi-Tenant Config (Isolated Revenue Streams)."
    ])

    # Slide 4: Real-World Example A - FitIt
    add_bullet_slide("Example A: FitIt (Fitness App)", [
        "Concept: Workout & Gym Discovery Platform.",
        "Key Microservices Used:",
        "- Media API: Heavy usage for Workout Video streaming.",
        "- Geo API: Location-based Gym finders.",
        "- Subscriptions: Monthly recurring billing."
    ])

    # Slide 5: Real-World Example B - Wissler
    add_bullet_slide("Example B: Wissler (Dating App)", [
        "Concept: Tinder-style matchmaking application.",
        "Key Microservices Used:",
        "- Geo API: Real-time 'People Nearby' queries.",
        "- Chat API: Heavy real-time messaging (MongoDB + SignalR).",
        "- Recommendation API: Matching algorithms.",
        "- Identity: Strict eKYC verification."
    ])

    # Slide 6: New Service Integration (eKYC)
    add_bullet_slide("Case Study: Adding eKYC", [
        "Scenario: Add Identity Verification for Fintech/Dating apps.",
        "Implementation:",
        "1. Create Verification.API microservice.",
        "2. Listen for 'UserRegistered' events via RabbitMQ.",
        "3. Store documents in secure 'VerificationDb'.",
        "4. Process via 3rd party (Onfido/SumSub)."
    ])

    # Slide 7: Tech Stack
    add_bullet_slide("Technology Stack", [
        "Backend: .NET 8, Entity Framework Core.",
        "Frontend: React 19, Tailwind CSS, Shadcn/UI.",
        "Deployment: Docker Compose (Dev), Kubernetes (Prod).",
        "Observability: Seq (Logs), HealthChecks UI.",
        "Gateway: YARP (Reverse Proxy) + MassTransit (Event Bus)."
    ])

    # Slide 8: Service URLs
    add_bullet_slide("Service URLs & Credentials", [
        "Global Admin: http://localhost:3000 (admin@globaldashboard.com / Password123!)",
        "FitIt Admin: http://localhost:3001 (admin@fitit.com / Password123!)",
        "Wissler Admin: http://localhost:3002 (admin@wissler.com / Password123!)",
        "Seq Logs: http://localhost:5341",
        "Auth API Swagger: http://localhost:5001/swagger"
    ])

    prs.save('Unified_Microservices_Platform.pptx')
    print("Presentation created successfully: Unified_Microservices_Platform.pptx")

if __name__ == "__main__":
    create_presentation()
